import os
import json
import tempfile
from typing import Dict, Any, List, Optional
import io
import base64
from ..config.supabase import supabase_client

# Importation conditionnelle des bibliothèques d'extraction
try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

try:
    import docx
    from docx.document import Document as DocxDocument
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import _Cell, Table
    from docx.text.paragraph import Paragraph
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.shapes.picture import Picture
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

async def extract_chapter_content(chapter_id: str) -> Dict[str, Any]:
    """
    Extrait le contenu d'un chapitre spécifique et le convertit en JSON
    
    Args:
        chapter_id: ID du chapitre à extraire
        
    Returns:
        Dictionnaire contenant le contenu extrait au format JSON
    """
    try:
        # 1. Récupérer les informations du chapitre depuis Supabase
        chapter_response = supabase_client.table("chapters").select("*").eq("id", chapter_id).execute()
        
        if not chapter_response.data or len(chapter_response.data) == 0:
            raise ValueError(f"Chapitre non trouvé: {chapter_id}")
        
        chapter = chapter_response.data[0]
        
        # 2. Récupérer le fichier depuis Supabase Storage
        file_path = chapter.get("file_path")
        if not file_path:
            raise ValueError(f"Chemin de fichier non trouvé pour le chapitre: {chapter_id}")
        
        # Extraire le bucket et le chemin du fichier
        parts = file_path.split("/")
        bucket_name = parts[0]  # Généralement 'chapters'
        file_name = "/".join(parts[1:])
        
        # Télécharger le fichier
        response = supabase_client.storage.from_(bucket_name).download(file_name)
        
        if not response:
            raise ValueError(f"Fichier non trouvé dans le stockage: {file_path}")
        
        # 3. Déterminer le type de fichier et extraire le contenu
        file_extension = os.path.splitext(file_name)[1].lower()
        
        # Récupérer le type de contenu (cours, exercice, examen)
        content_type = chapter.get("content_type", "course")
        
        # Créer un fichier temporaire pour traiter le document
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
            temp_file.write(response)
            temp_file_path = temp_file.name
        
        try:
            # Extraire le contenu selon le type de fichier
            if file_extension == '.pdf':
                content = extract_pdf_content(temp_file_path)
            elif file_extension == '.docx':
                content = extract_docx_content(temp_file_path)
            elif file_extension in ['.pptx', '.ppt']:
                content = extract_pptx_content(temp_file_path)
            else:
                content = {
                    "title": chapter.get("title", "Document inconnu"),
                    "type": "unknown",
                    "pages": [{"text": f"Format non pris en charge: {file_extension}"}],
                    "metadata": {"format": file_extension}
                }
            
            # Ajouter le type de contenu
            content["content_type"] = content_type
        finally:
            # Supprimer le fichier temporaire
            if os.path.exists(temp_file_path):
                os.unlink(temp_file_path)
        
        # 4. Mettre à jour le chapitre avec le contenu JSON extrait
        update_response = supabase_client.table("chapters").update(
            {"json_data": content}
        ).eq("id", chapter_id).execute()
        
        if update_response.error:
            raise ValueError(f"Erreur lors de la mise à jour du chapitre: {update_response.error}")
        
        return content
        
    except Exception as e:
        # Journaliser l'erreur et la renvoyer
        print(f"Erreur lors de l'extraction du chapitre {chapter_id}: {str(e)}")
        return {
            "title": f"Erreur - Chapitre {chapter_id}",
            "type": "error",
            "content_type": "course",
            "pages": [{"text": f"Erreur lors de l'extraction: {str(e)}"}],
            "metadata": {"error": str(e)}
        }

async def extract_all_course_chapters(course_id: str) -> int:
    """
    Extrait le contenu de tous les chapitres d'un cours
    
    Args:
        course_id: ID du cours dont les chapitres doivent être extraits
        
    Returns:
        Nombre de chapitres traités
    """
    try:
        # 1. Récupérer tous les chapitres du cours
        chapters_response = supabase_client.table("chapters").select("id").eq("course_id", course_id).execute()
        
        if not chapters_response.data:
            return 0
        
        # 2. Extraire le contenu de chaque chapitre
        processed_count = 0
        for chapter in chapters_response.data:
            try:
                await extract_chapter_content(chapter["id"])
                processed_count += 1
            except Exception as e:
                # Journaliser l'erreur mais continuer avec les autres chapitres
                print(f"Erreur lors de l'extraction du chapitre {chapter['id']}: {str(e)}")
        
        return processed_count
        
    except Exception as e:
        print(f"Erreur lors de l'extraction des chapitres du cours {course_id}: {str(e)}")
        raise e

def extract_pdf_content(file_path: str) -> Dict[str, Any]:
    """
    Extrait le contenu d'un fichier PDF
    
    Args:
        file_path: Chemin vers le fichier PDF
        
    Returns:
        Dictionnaire contenant le contenu extrait
    """
    if not PYMUPDF_AVAILABLE:
        return {
            "title": os.path.basename(file_path),
            "type": "pdf",
            "pages": [{"text": "PyMuPDF non disponible pour l'extraction PDF"}],
            "metadata": {"error": "PyMuPDF non installé"}
        }
    
    try:
        # Ouvrir le document PDF
        doc = fitz.open(file_path)
        
        # Extraire les métadonnées
        metadata = {
            "title": doc.metadata.get("title", ""),
            "author": doc.metadata.get("author", ""),
            "subject": doc.metadata.get("subject", ""),
            "keywords": doc.metadata.get("keywords", ""),
            "page_count": len(doc)
        }
        
        # Extraire le contenu de chaque page
        pages = []
        for page_num, page in enumerate(doc):
            # Extraire le texte
            text = page.get_text()
            
            # Extraire les images (optionnel, peut être lourd)
            images = []
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                if base_image:
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                    images.append({
                        "index": img_index,
                        "format": image_ext,
                        "data": f"data:image/{image_ext};base64,{image_base64}"
                    })
            
            pages.append({
                "page_num": page_num + 1,
                "text": text,
                "images": images
            })
        
        # Fermer le document
        doc.close()
        
        return {
            "title": metadata.get("title") or os.path.basename(file_path),
            "type": "pdf",
            "pages": pages,
            "metadata": metadata
        }
        
    except Exception as e:
        print(f"Erreur lors de l'extraction du PDF {file_path}: {str(e)}")
        return {
            "title": os.path.basename(file_path),
            "type": "pdf",
            "pages": [{"text": f"Erreur lors de l'extraction: {str(e)}"}],
            "metadata": {"error": str(e)}
        }

def extract_docx_content(file_path: str) -> Dict[str, Any]:
    """
    Extrait le contenu d'un fichier DOCX
    
    Args:
        file_path: Chemin vers le fichier DOCX
        
    Returns:
        Dictionnaire contenant le contenu extrait
    """
    if not PYTHON_DOCX_AVAILABLE or not PIL_AVAILABLE:
        return {
            "title": os.path.basename(file_path),
            "type": "docx",
            "pages": [{"text": "python-docx ou PIL non disponible pour l'extraction DOCX"}],
            "metadata": {"error": "python-docx ou PIL non installé"}
        }
    
    try:
        # Ouvrir le document DOCX
        doc = docx.Document(file_path)
        
        # Extraire les métadonnées
        core_properties = doc.core_properties
        metadata = {
            "title": core_properties.title or "",
            "author": core_properties.author or "",
            "subject": core_properties.subject or "",
            "keywords": core_properties.keywords or "",
            "paragraph_count": len(doc.paragraphs),
            "section_count": len(doc.sections)
        }
        
        # Fonction pour extraire le contenu structuré
        def iter_block_items(parent):
            """
            Parcourt tous les éléments du document (paragraphes et tables)
            """
            if isinstance(parent, DocxDocument):
                parent_elem = parent.element.body
            elif isinstance(parent, _Cell):
                parent_elem = parent._tc
            else:
                raise ValueError("Élément parent non pris en charge")

            for child in parent_elem.iterchildren():
                if isinstance(child, CT_P):
                    yield Paragraph(child, parent)
                elif isinstance(child, CT_Tbl):
                    yield Table(child, parent)
        
        # Extraire le contenu structuré
        content_blocks = []
        for block in iter_block_items(doc):
            if isinstance(block, Paragraph):
                if block.text.strip():
                    # Extraire le style et la mise en forme
                    style_name = block.style.name if block.style else "Normal"
                    is_heading = style_name.startswith('Heading')
                    heading_level = int(style_name.replace('Heading ', '')) if is_heading else 0
                    
                    # Extraire le texte avec sa mise en forme
                    formatted_text = []
                    for run in block.runs:
                        formatted_text.append({
                            "text": run.text,
                            "bold": run.bold,
                            "italic": run.italic,
                            "underline": run.underline,
                            "font": run.font.name if run.font.name else "default",
                            "size": run.font.size if run.font.size else "default",
                            "color": run.font.color.rgb if run.font.color and run.font.color.rgb else None
                        })
                    
                    content_blocks.append({
                        "type": "paragraph",
                        "text": block.text,
                        "formatted_text": formatted_text,
                        "style": style_name,
                        "is_heading": is_heading,
                        "heading_level": heading_level,
                        "alignment": str(block.alignment) if block.alignment else "left"
                    })
            elif isinstance(block, Table):
                # Extraire les données du tableau
                table_data = []
                for i, row in enumerate(block.rows):
                    row_data = []
                    for j, cell in enumerate(row.cells):
                        cell_text = cell.text.strip()
                        row_data.append(cell_text)
                    table_data.append(row_data)
                
                content_blocks.append({
                    "type": "table",
                    "rows": len(block.rows),
                    "columns": len(block.columns),
                    "data": table_data
                })
        
        # Extraire les images
        images = []
        for rel in doc.part.rels.values():
            if rel.reltype == docx.opc.constants.RELATIONSHIP_TYPE.IMAGE:
                try:
                    image_bytes = rel.target_part.blob
                    image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                    image_format = rel.target_ref.split('.')[-1].lower()
                    
                    images.append({
                        "id": rel.rId,
                        "format": image_format,
                        "data": f"data:image/{image_format};base64,{image_base64}"
                    })
                except Exception as img_err:
                    print(f"Erreur lors de l'extraction d'une image: {str(img_err)}")
        
        # Simuler des pages (DOCX n'a pas de concept de pages)
        pages = [{
            "page_num": 1, 
            "content_blocks": content_blocks,
            "images": images,
            "has_images": len(images) > 0,
            "has_tables": any(block.get("type") == "table" for block in content_blocks)
        }]
        
        return {
            "title": metadata.get("title") or os.path.basename(file_path),
            "type": "docx",
            "pages": pages,
            "metadata": metadata,
            "structure": {
                "has_images": len(images) > 0,
                "has_tables": any(block.get("type") == "table" for block in content_blocks),
                "heading_count": sum(1 for block in content_blocks if block.get("is_heading", False))
            }
        }
        
    except Exception as e:
        print(f"Erreur lors de l'extraction du DOCX {file_path}: {str(e)}")
        return {
            "title": os.path.basename(file_path),
            "type": "docx",
            "pages": [{"text": f"Erreur lors de l'extraction: {str(e)}"}],
            "metadata": {"error": str(e)}
        }
        
    except Exception as e:
        print(f"Erreur lors de l'extraction du DOCX {file_path}: {str(e)}")
        return {
            "title": os.path.basename(file_path),
            "type": "docx",
            "pages": [{"text": f"Erreur lors de l'extraction: {str(e)}"}],
            "metadata": {"error": str(e)}
        }

def extract_pptx_content(file_path: str) -> Dict[str, Any]:
    """
    Extrait le contenu d'un fichier PPTX
    
    Args:
        file_path: Chemin vers le fichier PPTX
        
    Returns:
        Dictionnaire contenant le contenu extrait
    """
    if not PYTHON_PPTX_AVAILABLE or not PIL_AVAILABLE:
        return {
            "title": os.path.basename(file_path),
            "type": "pptx",
            "pages": [{"text": "python-pptx ou PIL non disponible pour l'extraction PPTX"}],
            "metadata": {"error": "python-pptx ou PIL non installé"}
        }
    
    try:
        # Ouvrir la présentation
        prs = Presentation(file_path)
        
        # Extraire les métadonnées
        core_properties = prs.core_properties
        metadata = {
            "slide_count": len(prs.slides),
            "title": core_properties.title or "",
            "author": core_properties.author or "",
            "subject": core_properties.subject or "",
            "keywords": core_properties.keywords or ""
        }
        
        # Extraire le contenu de chaque diapositive
        slides = []
        for slide_num, slide in enumerate(prs.slides):
            # Extraire les textes avec structure
            content_blocks = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Déterminer le type de texte (titre, sous-titre, contenu)
                    shape_type = "content"
                    if hasattr(shape, "is_title") and shape.is_title:
                        shape_type = "title"
                    elif hasattr(shape, "name") and "Title" in shape.name:
                        shape_type = "title"
                    elif hasattr(shape, "name") and "Subtitle" in shape.name:
                        shape_type = "subtitle"
                    
                    # Extraire le texte avec mise en forme
                    formatted_text = []
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            formatted_text.append({
                                "text": run.text,
                                "bold": run.font.bold,
                                "italic": run.font.italic,
                                "underline": run.font.underline,
                                "font": run.font.name if run.font.name else "default",
                                "size": run.font.size if run.font.size else "default",
                                "color": str(run.font.color.rgb) if run.font.color and run.font.color.rgb else None
                            })
                    
                    content_blocks.append({
                        "type": "text",
                        "text_type": shape_type,
                        "text": shape.text,
                        "formatted_text": formatted_text,
                        "position": {
                            "x": shape.left,
                            "y": shape.top,
                            "width": shape.width,
                            "height": shape.height
                        }
                    })
                elif isinstance(shape, Picture):
                    # Ajouter l'information sur l'image (sans l'extraire encore)
                    content_blocks.append({
                        "type": "image_placeholder",
                        "position": {
                            "x": shape.left,
                            "y": shape.top,
                            "width": shape.width,
                            "height": shape.height
                        }
                    })
                elif hasattr(shape, "has_table") and shape.has_table:
                    # Extraire les données du tableau
                    table_data = []
                    for row in shape.table.rows:
                        row_data = []
                        for cell in row.cells:
                            if cell.text_frame:
                                row_data.append(cell.text_frame.text)
                            else:
                                row_data.append("")
                        table_data.append(row_data)
                    
                    content_blocks.append({
                        "type": "table",
                        "rows": len(shape.table.rows),
                        "columns": len(shape.table.columns),
                        "data": table_data,
                        "position": {
                            "x": shape.left,
                            "y": shape.top,
                            "width": shape.width,
                            "height": shape.height
                        }
                    })
            
            # Extraire les images
            images = []
            for shape in slide.shapes:
                if isinstance(shape, Picture):
                    try:
                        # Extraire l'image
                        image_bytes = shape.image.blob
                        image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                        image_type = shape.image.ext.lower()
                        
                        images.append({
                            "format": image_type,
                            "data": f"data:image/{image_type};base64,{image_base64}",
                            "position": {
                                "x": shape.left,
                                "y": shape.top,
                                "width": shape.width,
                                "height": shape.height
                            }
                        })
                    except Exception as img_err:
                        print(f"Erreur lors de l'extraction d'une image: {str(img_err)}")
            
            # Créer la représentation de la diapositive
            slide_text = "\n".join([block["text"] for block in content_blocks if block["type"] == "text"])
            slides.append({
                "slide_num": slide_num + 1,
                "text": slide_text,  # Pour compatibilité avec l'ancienne version
                "content_blocks": content_blocks,
                "images": images,
                "has_images": len(images) > 0,
                "has_tables": any(block.get("type") == "table" for block in content_blocks)
            })
        
        return {
            "title": metadata.get("title") or os.path.basename(file_path),
            "type": "pptx",
            "pages": slides,
            "metadata": metadata,
            "structure": {
                "total_slides": len(slides),
                "has_images": any(slide.get("has_images", False) for slide in slides),
                "has_tables": any(slide.get("has_tables", False) for slide in slides)
            }
        }
        
    except Exception as e:
        print(f"Erreur lors de l'extraction du PPTX {file_path}: {str(e)}")
        return {
            "title": os.path.basename(file_path),
            "type": "pptx",
            "pages": [{"text": f"Erreur lors de l'extraction: {str(e)}"}],
            "metadata": {"error": str(e)}
        }
